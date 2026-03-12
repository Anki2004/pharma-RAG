"""
app.py
------
FastAPI application with:
  - GET /search?query=<text>        → Top-5 RAG results + AI summary
  - GET /export/excel?query=<text>  → Download .xlsx
  - GET /export/pdf?query=<text>    → Download .pdf
  - GET /export/ppt?query=<text>    → Download .pptx
  - GET /export/graph?query=<text>  → Download .png bar chart
"""

import logging
import io
import json
import requests
import pandas as pd
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt

# Suppress HuggingFace transformers warnings (position_ids, etc.)
logging.getLogger("transformers").setLevel(logging.ERROR)

from fastapi import FastAPI, Query
from fastapi.responses import JSONResponse, StreamingResponse, HTMLResponse
from sentence_transformers import SentenceTransformer
from pinecone import Pinecone

from openpyxl import Workbook
from fpdf import FPDF
from pptx import Presentation
from pptx.util import Inches, Pt

from config import (
    PINECONE_API_KEY,
    PINECONE_INDEX_NAME,
    EMBEDDING_MODEL,
    OLLAMA_BASE_URL,
    OLLAMA_MODEL,
    TOP_K,
)

# ── Startup ──────────────────────────────────────────────────────────────

app = FastAPI(title="Pharma RAG Search", version="1.0.0")

# Load embedding model once
embed_model = SentenceTransformer(EMBEDDING_MODEL)

# Connect to Pinecone
pc = Pinecone(api_key=PINECONE_API_KEY)
index = pc.Index(PINECONE_INDEX_NAME)


# ── Helpers ──────────────────────────────────────────────────────────────

def search_pinecone(query: str, top_k: int = TOP_K) -> list[dict]:
    """Embed the query and search Pinecone. Returns list of results."""
    query_embedding = embed_model.encode(query).tolist()
    response = index.query(vector=query_embedding, top_k=top_k, include_metadata=True)

    results = []
    for match in response.matches:
        results.append({
            "id": match.id,
            "score": round(match.score, 4),
            "title": match.metadata.get("title", ""),
            "description": match.metadata.get("description", ""),
            "source": match.metadata.get("source", ""),
            "industry": match.metadata.get("industry", ""),
        })
    return results


def generate_summary(query: str, results: list[dict]) -> str:
    """Call Ollama local Mistral model to summarize the retrieved results."""
    context = "\n".join(
        [f"- {r['description']} (source: {r['source']})" for r in results]
    )
    prompt = f"""You are a helpful assistant for the chemical and pharmaceutical industry.
Based on the following search results, provide a clear and concise summary answering the user's query.

User Query: {query}

Retrieved Data:
{context}

Provide a helpful summary:"""

    try:
        resp = requests.post(
            f"{OLLAMA_BASE_URL}/api/generate",
            json={
                "model": OLLAMA_MODEL,
                "prompt": prompt,
                "stream": False,
            },
            timeout=60,
        )
        resp.raise_for_status()
        return resp.json().get("response", "No summary generated.")
    except Exception as e:
        return f"AI summary unavailable (Ollama error: {e})"


# ── Endpoints ────────────────────────────────────────────────────────────

@app.get("/", response_class=HTMLResponse)
def root():
    return """
    <!DOCTYPE html>
    <html>
    <head>
        <title>Pharma RAG Search</title>
        <style>
            body { font-family: sans-serif; max-width: 800px; margin: 0 auto; padding: 20px; }
            .search-box { margin: 20px 0; display: flex; gap: 10px; }
            input { flex: 1; padding: 10px; font-size: 16px; }
            button { padding: 10px 20px; font-size: 16px; cursor: pointer; background: #007bff; color: white; border: none; }
            button:hover { background: #0056b3; }
            .results { margin-top: 20px; }
            .result-item { border: 1px solid #ddd; padding: 15px; margin-bottom: 10px; border-radius: 5px; }
            .downloads { margin-top: 20px; }
            .downloads a { margin-right: 15px; text-decoration: none; color: #007bff; }
            .summary { background: #f0f8ff; padding: 15px; border-radius: 5px; margin-bottom: 20px; }
        </style>
    </head>
    <body>
        <h1>Pharma Product Knowledge Search</h1>
        <div class="search-box">
            <input type="text" id="query" placeholder="e.g., Suppliers of Metformin in Europe" />
            <button onclick="search()">Search</button>
        </div>
        
        <div id="loading" style="display:none;">Searching...</div>
        
        <div id="results-area" style="display:none;">
            <div class="summary">
                <h3>AI Summary</h3>
                <p id="ai-summary"></p>
            </div>
            
            <div class="downloads">
                <strong>Export Results:</strong>
                <a id="link-excel" href="#" target="_blank">Excel</a>
                <a id="link-pdf" href="#" target="_blank">PDF</a>
                <a id="link-ppt" href="#" target="_blank">PowerPoint</a>
                <a id="link-graph" href="#" target="_blank">Graph</a>
            </div>
            
            <h3>Detailed Results</h3>
            <div id="results-list"></div>
        </div>

        <script>
            async function search() {
                const query = document.getElementById('query').value;
                if (!query) return;
                
                document.getElementById('loading').style.display = 'block';
                document.getElementById('results-area').style.display = 'none';
                
                try {
                    const res = await fetch(`/search?query=${encodeURIComponent(query)}`);
                    const data = await res.json();
                    
                    document.getElementById('ai-summary').textContent = data.ai_summary;
                    
                    const list = document.getElementById('results-list');
                    list.innerHTML = '';
                    data.results.forEach(r => {
                        const div = document.createElement('div');
                        div.className = 'result-item';
                        div.innerHTML = `<strong>${r.title}</strong> (Score: ${r.score})<br>${r.description}<br><small>Source: ${r.source}</small>`;
                        list.appendChild(div);
                    });
                    
                    // Update export links
                    document.getElementById('link-excel').href = `/export/excel?query=${encodeURIComponent(query)}`;
                    document.getElementById('link-pdf').href = `/export/pdf?query=${encodeURIComponent(query)}`;
                    document.getElementById('link-ppt').href = `/export/ppt?query=${encodeURIComponent(query)}`;
                    document.getElementById('link-graph').href = `/export/graph?query=${encodeURIComponent(query)}`;
                    
                    document.getElementById('results-area').style.display = 'block';
                } catch (e) {
                    alert("Search failed: " + e);
                } finally {
                    document.getElementById('loading').style.display = 'none';
                }
            }
        </script>
    </body>
    </html>
    """


@app.get("/search")
def search(query: str = Query(..., description="Natural language search query")):
    """Search the knowledge base and return Top-K results with AI summary."""
    results = search_pinecone(query)
    ai_summary = generate_summary(query, results)
    return JSONResponse({
        "query": query,
        "results": results,
        "ai_summary": ai_summary,
    })


# ── Export: Excel ────────────────────────────────────────────────────────

@app.get("/export/excel")
def export_excel(query: str = Query(..., description="Search query")):
    """Export search results as an Excel file."""
    results = search_pinecone(query)

    wb = Workbook()
    ws = wb.active
    ws.title = "Search Results"
    ws.append(["ID", "Score", "Title", "Description", "Source"])
    for r in results:
        ws.append([r["id"], r["score"], r["title"], r["description"], r["source"]])

    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        headers={"Content-Disposition": f'attachment; filename="results.xlsx"'},
    )


# ── Export: PDF ──────────────────────────────────────────────────────────

@app.get("/export/pdf")
def export_pdf(query: str = Query(..., description="Search query")):
    """Export search results as a PDF file."""
    results = search_pinecone(query)

    pdf = FPDF()
    pdf.add_page()
    pdf.set_font("Helvetica", "B", 16)
    pdf.cell(0, 10, "Pharma RAG Search Results", ln=True, align="C")
    pdf.ln(5)
    pdf.set_font("Helvetica", "", 11)
    pdf.cell(0, 8, f"Query: {query}", ln=True)
    pdf.ln(5)

    for i, r in enumerate(results, 1):
        pdf.set_font("Helvetica", "B", 12)
        pdf.cell(0, 8, f"{i}. {r['title']}  (score: {r['score']})", ln=True)
        pdf.set_font("Helvetica", "", 10)
        # Wrap long descriptions
        pdf.multi_cell(0, 6, f"   {r['description']}")
        pdf.cell(0, 6, f"   Source: {r['source']}", ln=True)
        pdf.ln(3)

    buf = io.BytesIO(pdf.output())
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/pdf",
        headers={"Content-Disposition": f'attachment; filename="results.pdf"'},
    )


# ── Export: PowerPoint ───────────────────────────────────────────────────

@app.get("/export/ppt")
def export_ppt(query: str = Query(..., description="Search query")):
    """Export search results as a PowerPoint file."""
    results = search_pinecone(query)

    prs = Presentation()
    # Title slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Pharma RAG Search Results"
    slide.placeholders[1].text = f"Query: {query}"

    # One slide per result
    for r in results:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = r["title"]
        body = slide.placeholders[1]
        body.text = f"Score: {r['score']}\n\n{r['description']}\n\nSource: {r['source']}"

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f'attachment; filename="results.pptx"'},
    )


# ── Export: Graph (PNG) ──────────────────────────────────────────────────

@app.get("/export/graph")
def export_graph(query: str = Query(..., description="Search query")):
    """Export a bar chart of search result scores as PNG."""
    results = search_pinecone(query)

    titles = [r["title"][:20] for r in results]
    scores = [r["score"] for r in results]

    fig, ax = plt.subplots(figsize=(8, 5))
    bars = ax.barh(titles, scores, color="#4f8ef7")
    ax.set_xlabel("Similarity Score")
    ax.set_title(f"Search Results: {query}")
    ax.invert_yaxis()
    plt.tight_layout()

    buf = io.BytesIO()
    fig.savefig(buf, format="png", dpi=150)
    plt.close(fig)
    buf.seek(0)

    return StreamingResponse(
        buf,
        media_type="image/png",
        headers={"Content-Disposition": f'attachment; filename="results.png"'},
    )


# ── Run ──────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    import uvicorn
    uvicorn.run("app:app", host="0.0.0.0", port=8000, reload=True)
